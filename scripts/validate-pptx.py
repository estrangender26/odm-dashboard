#!/usr/bin/env python3
"""
PPTX Validation Script for Governance Presentation
Usage: python3 validate-pptx.py <pptx_path>
"""

import sys
import json

try:
    from pptx import Presentation
    
    pptx_path = sys.argv[1] if len(sys.argv) > 1 else 'validation-artifacts/governance-final-validation.pptx'
    prs = Presentation(pptx_path)
    
    results = {
        'total_slides': len(prs.slides), 
        'slides': []
    }
    
    for idx, slide in enumerate(prs.slides, 1):
        slide_data = {'slide_number': idx, 'text': []}
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                slide_data['text'].append(shape.text.strip())
        results['slides'].append(slide_data)
    
    print(json.dumps(results))
    sys.exit(0)
    
except Exception as e:
    print(json.dumps({'error': str(e), 'total_slides': 0, 'slides': []}))
    sys.exit(1)
